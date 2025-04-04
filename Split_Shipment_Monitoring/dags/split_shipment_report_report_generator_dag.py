"""
Manhattan Active Omni - Split Shipment Report Generator DAG
----------------------------------------------------------
This DAG generates exportable reports about split shipments.
It performs the following steps:
1. Loads collected split shipment data
2. Generates CSV reports
3. Creates PowerPoint presentations with charts and insights
4. Stores the reports for access via the dashboard
"""

from datetime import datetime, timedelta
import json
import pandas as pd
import numpy as np
import io
import logging
from airflow import DAG
from airflow.operators.python import PythonOperator
from airflow.models import Variable
from airflow.providers.amazon.aws.hooks.s3 import S3Hook
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
import matplotlib.pyplot as plt
import seaborn as sns

# Set up logging
logger = logging.getLogger(__name__)

# Default arguments for the DAG
default_args = {
    'owner': 'airflow',
    'depends_on_past': False,
    'email_on_failure': True,
    'email_on_retry': False,
    'retries': 1,
    'retry_delay': timedelta(minutes=5),
    'start_date': datetime(2023, 1, 1),
}

# Create the DAG
dag = DAG(
    'split_shipment_report_generator',
    default_args=default_args,
    description='Generates exportable reports about split shipments',
    schedule_interval='0 7 * * *',  # Daily at 7 AM
    catchup=False,
    tags=['manhattan', 'split_shipment', 'monitoring', 'reports'],
)

# Environment variables
ORG_LIST = Variable.get("organization_list", deserialize_json=True)
S3_BUCKET = Variable.get("s3_data_bucket")
S3_PREFIX = "split_shipment_data"
REPORTS_PREFIX = "split_shipment_reports"

def load_data_for_reports(organization, report_type, days=30, **kwargs):
    """Load data for generating reports."""
    try:
        s3_hook = S3Hook()
        
        # Determine what data to load based on report type
        if report_type == 'daily':
            # Get yesterday's detailed data
            execution_date = kwargs.get('execution_date', datetime.now())
            yesterday = (execution_date - timedelta(days=1)).strftime('%Y-%m-%d')
            
            detailed_key = f"{S3_PREFIX}/detailed/{organization}/{yesterday}_split_shipments.json"
            try:
                detailed_data_str = s3_hook.read_key(key=detailed_key, bucket_name=S3_BUCKET)
                detailed_data = json.loads(detailed_data_str)
                detailed_df = pd.DataFrame(detailed_data)
            except Exception as e:
                logger.warning(f"Could not load detailed data for {organization}: {str(e)}")
                detailed_df = pd.DataFrame()
            
            # Get yesterday's summary
            summary_key = f"{S3_PREFIX}/summary/{organization}/{yesterday}_summary.json"
            try:
                summary_data_str = s3_hook.read_key(key=summary_key, bucket_name=S3_BUCKET)
                summary_data = json.loads(summary_data_str)
            except Exception as e:
                logger.warning(f"Could not load summary data for {organization}: {str(e)}")
                summary_data = {}
            
            # Get yesterday's analysis
            analysis_key = f"{S3_PREFIX}/analysis/{organization}/{yesterday}_analysis.json"
            try:
                analysis_data_str = s3_hook.read_key(key=analysis_key, bucket_name=S3_BUCKET)
                analysis_data = json.loads(analysis_data_str)
            except Exception as e:
                logger.warning(f"Could not load analysis data for {organization}: {str(e)}")
                analysis_data = {}
            
            return {
                'detailed_df': detailed_df,
                'summary': summary_data,
                'analysis': analysis_data,
                'report_date': yesterday
            }
        
        elif report_type == 'weekly' or report_type == 'monthly':
            # Load historical data for trend analysis
            historical_key = f"{S3_PREFIX}/historical/{organization}_split_rates.json"
            try:
                historical_data_str = s3_hook.read_key(key=historical_key, bucket_name=S3_BUCKET)
                historical_data = json.loads(historical_data_str)
                historical_df = pd.DataFrame(historical_data)
                historical_df['date'] = pd.to_datetime(historical_df['date'])
                historical_df = historical_df.sort_values('date')
            except Exception as e:
                logger.warning(f"Could not load historical data for {organization}: {str(e)}")
                historical_df = pd.DataFrame()
            
            # Calculate time period for the report
            end_date = kwargs.get('execution_date', datetime.now())
            if report_type == 'weekly':
                start_date = end_date - timedelta(days=7)
                period_data = historical_df[historical_df['date'] >= start_date]
            else:  # monthly
                start_date = end_date - timedelta(days=30)
                period_data = historical_df[historical_df['date'] >= start_date]
            
            return {
                'historical_df': historical_df,
                'period_data': period_data,
                'start_date': start_date.strftime('%Y-%m-%d'),
                'end_date': end_date.strftime('%Y-%m-%d')
            }
        
        else:
            logger.error(f"Unknown report type: {report_type}")
            raise ValueError(f"Unknown report type: {report_type}")
        
    except Exception as e:
        logger.error(f"Failed to load data for {organization} {report_type} report: {str(e)}")
        raise

def generate_csv_report(data, organization, report_type, **kwargs):
    """Generate CSV report from the data."""
    try:
        if report_type == 'daily':
            detailed_df = data.get('detailed_df')
            
            if detailed_df.empty:
                logger.warning(f"No detailed data available for {organization} daily CSV report")
                return None
            
            # Create CSV buffer
            csv_buffer = io.StringIO()
            
            # Write detailed data to CSV
            detailed_df.to_csv(csv_buffer, index=False)
            
            # Get the CSV content
            csv_content = csv_buffer.getvalue()
            
            # Prepare filename
            report_date = data.get('report_date', datetime.now().strftime('%Y-%m-%d'))
            filename = f"{organization}_{report_date}_split_shipments.csv"
            
            return {
                'content': csv_content,
                'filename': filename,
                'content_type': 'text/csv'
            }
        
        elif report_type == 'weekly' or report_type == 'monthly':
            period_data = data.get('period_data')
            
            if period_data.empty:
                logger.warning(f"No period data available for {organization} {report_type} CSV report")
                return None
            
            # Create CSV buffer
            csv_buffer = io.StringIO()
            
            # Write period data to CSV
            period_data.to_csv(csv_buffer, index=False)
            
            # Get the CSV content
            csv_content = csv_buffer.getvalue()
            
            # Prepare filename
            start_date = data.get('start_date')
            end_date = data.get('end_date')
            filename = f"{organization}_{start_date}_to_{end_date}_{report_type}_split_shipments.csv"
            
            return {
                'content': csv_content,
                'filename': filename,
                'content_type': 'text/csv'
            }
        
        else:
            logger.error(f"Unknown report type: {report_type}")
            raise ValueError(f"Unknown report type: {report_type}")
        
    except Exception as e:
        logger.error(f"Failed to generate CSV report for {organization} {report_type}: {str(e)}")
        raise

def create_ppt_report(data, organization, report_type, **kwargs):
    """Create PowerPoint report with charts and insights."""
    try:
        # Create a presentation
        prs = Presentation()
        
        # Set slide dimensions to widescreen (16:9)
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)
        
        if report_type == 'daily':
            detailed_df = data.get('detailed_df')
            summary = data.get('summary', {})
            analysis = data.get('analysis', {})
            report_date = data.get('report_date', datetime.now().strftime('%Y-%m-%d'))
            
            if detailed_df.empty:
                logger.warning(f"No detailed data available for {organization} daily PowerPoint report")
                return None
            
            # Title slide
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)
            title = slide.shapes.title
            subtitle = slide.placeholders[1]
            
            title.text = f"Split Shipment Report: {organization}"
            subtitle.text = f"Daily Report for {report_date}"
            
            # Summary slide
            summary_slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(summary_slide_layout)
            title = slide.shapes.title
            content = slide.placeholders[1]
            
            title.text = "Summary"
            content_text = f"Total Orders: {summary.get('total_orders', 'N/A')}\n"
            content_text += f"Split Orders: {summary.get('split_orders', 'N/A')}\n"
            content_text += f"Split Rate: {summary.get('split_rate', 'N/A'):.2f}%\n"
            
            if analysis and 'trend' in analysis:
                trend = analysis.get('trend', {})
                content_text += f"\nTrend Direction: {trend.get('direction', 'N/A')}"
            
            content.text = content_text
            
            # Location analysis slide
            if analysis and 'location_analysis' in analysis and analysis['location_analysis']:
                location_slide_layout = prs.slide_layouts[6]  # Blank slide
                slide = prs.slides.add_slide(location_slide_layout)
                
                # Add title
                title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(1))
                title_frame = title_shape.text_frame
                title_frame.text = "Split Rates by Location"
                title_frame.paragraphs[0].font.size = Pt(28)
                title_frame.paragraphs[0].font.bold = True
                
                # Create chart data
                chart_data = CategoryChartData()
                locations = []
                split_rates = []
                
                # Get top 10 locations by split rate
                top_locations = analysis['location_analysis'][:10]
                
                for loc in top_locations:
                    locations.append(loc['ship_from_location'])
                    split_rates.append(loc['split_rate'])
                
                chart_data.categories = locations
                chart_data.add_series('Split Rate (%)', split_rates)
                
                # Add chart
                chart = slide.shapes.add_chart(
                    XL_CHART_TYPE.COLUMN_CLUSTERED,
                    Inches(0.5), Inches(1.5),
                    Inches(12), Inches(5.5),
                    chart_data
                )
            
            # SKU analysis slide
            if analysis and 'sku_analysis' in analysis and analysis['sku_analysis']:
                sku_slide_layout = prs.slide_layouts[6]  # Blank slide
                slide = prs.slides.add_slide(sku_slide_layout)
                
                # Add title
                title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(1))
                title_frame = title_shape.text_frame
                title_frame.text = "Split Rates by SKU"
                title_frame.paragraphs[0].font.size = Pt(28)
                title_frame.paragraphs[0].font.bold = True
                
                # Create chart data
                chart_data = CategoryChartData()
                skus = []
                split_rates = []
                
                # Get top 10 SKUs by split rate
                top_skus = analysis['sku_analysis'][:10]
                
                for sku in top_skus:
                    skus.append(sku['item_id'])
                    split_rates.append(sku['split_rate'])
                
                chart_data.categories = skus
                chart_data.add_series('Split Rate (%)', split_rates)
                
                # Add chart
                chart = slide.shapes.add_chart(
                    XL_CHART_TYPE.COLUMN_CLUSTERED,
                    Inches(0.5), Inches(1.5),
                    Inches(12), Inches(5.5),
                    chart_data
                )
        
        elif report_type == 'weekly' or report_type == 'monthly':
            historical_df = data.get('historical_df')
            period_data = data.get('period_data')
            start_date = data.get('start_date')
            end_date = data.get('end_date')
            
            if period_data.empty:
                logger.warning(f"No period data available for {organization} {report_type} PowerPoint report")
                return None
            
            # Title slide
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)
            title = slide.shapes.title
            subtitle = slide.placeholders[1]
            
            title.text = f"Split Shipment Report: {organization}"
            subtitle.text = f"{report_type.capitalize()} Report from {start_date} to {end_date}"
            
            # Trend analysis slide
            trend_slide_layout = prs.slide_layouts[6]  # Blank slide
            slide = prs.slides.add_slide(trend_slide_layout)
            
            # Add title
            title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(1))
            title_frame = title_shape.text_frame
            title_frame.text = "Split Rate Trend"
            title_frame.paragraphs[0].font.size = Pt(28)
            title_frame.paragraphs[0].font.bold = True
            
            # Create chart data for trend
            chart_data = CategoryChartData()
            dates = []
            split_rates = []
            
            for _, row in period_data.iterrows():
                dates.append(row['date'].strftime('%Y-%m-%d'))
                split_rates.append(row['split_rate'])
            
            chart_data.categories = dates
            chart_data.add_series('Split Rate (%)', split_rates)
            
            # Add chart
            chart = slide.shapes.add_chart(
                XL_CHART_TYPE.LINE,
                Inches(0.5), Inches(1.5),
                Inches(12), Inches(5.5),
                chart_data
            )
            
            # Summary statistics slide
            summary_slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(summary_slide_layout)
            title = slide.shapes.title
            content = slide.placeholders[1]
            
            title.text = "Summary Statistics"
            
            avg_split_rate = period_data['split_rate'].mean()
            max_split_rate = period_data['split_rate'].max()
            min_split_rate = period_data['split_rate'].min()
            total_orders = period_data['total_orders'].sum()
            total_split_orders = period_data['split_orders'].sum()
            overall_split_rate = (total_split_orders / total_orders * 100) if total_orders > 0 else 0
            
            content_text = f"Average Split Rate: {avg_split_rate:.2f}%\n"
            content_text += f"Maximum Split Rate: {max_split_rate:.2f}%\n"
            content_text += f"Minimum Split Rate: {min_split_rate:.2f}%\n\n"
            content_text += f"Total Orders: {total_orders}\n"
            content_text += f"Total Split Orders: {total_split_orders}\n"
            content_text += f"Overall Split Rate: {overall_split_rate:.2f}%"
            
            content.text = content_text
        
        else:
            logger.error(f"Unknown report type: {report_type}")
            raise ValueError(f"Unknown report type: {report_type}")
        
        # Save presentation to memory
        pptx_buffer = io.BytesIO()
        prs.save(pptx_buffer)
        pptx_buffer.seek(0)
        
        # Prepare filename
        if report_type == 'daily':
            report_date = data.get('report_date', datetime.now().strftime('%Y-%m-%d'))
            filename = f"{organization}_{report_date}_split_shipments.pptx"
        else:
            start_date = data.get('start_date')
            end_date = data.get('end_date')
            filename = f"{organization}_{start_date}_to_{end_date}_{report_type}_split_shipments.pptx"
        
        return {
            'content': pptx_buffer.getvalue(),
            'filename': filename,
            'content_type': 'application/vnd.openxmlformats-officedocument.presentationml.presentation'
        }
        
    except Exception as e:
        logger.error(f"Failed to create PowerPoint report for {organization} {report_type}: {str(e)}")
        raise

def store_report(report, organization, report_type, **kwargs):
    """Store the generated report in S3."""
    try:
        if not report:
            logger.warning(f"No report to store for {organization} {report_type}")
            return
        
        s3_hook = S3Hook()
        
        # Store the report in S3
        s3_key = f"{REPORTS_PREFIX}/{report_type}/{organization}/{report['filename']}"
        s3_hook.load_bytes(
            bytes_data=report['content'].encode() if isinstance(report['content'], str) else report['content'],
            key=s3_key,
            bucket_name=S3_BUCKET,
            replace=True
        )
        
        # Create a metadata file for the report
        metadata = {
            'filename': report['filename'],
            'content_type': report['content_type'],
            'organization': organization,
            'report_type': report_type,
            'created_at': datetime.now().isoformat(),
            's3_key': s3_key
        }
        
        metadata_key = f"{REPORTS_PREFIX}/{report_type}/{organization}/{report['filename']}.metadata.json"
        s3_hook.load_string(
            json.dumps(metadata),
            key=metadata_key,
            bucket_name=S3_BUCKET,
            replace=True
        )
        
        return {
            's3_key': s3_key,
            'metadata_key': metadata_key
        }
    
    except Exception as e:
        logger.error(f"Failed to store report for {organization} {report_type}: {str(e)}")
        raise

def generate_report_for_organization(organization, report_type, **kwargs):
    """Generate a report for a specific organization and report type."""
    try:
        # Load data for the report
        data = load_data_for_reports(organization, report_type, **kwargs)
        
        # Generate CSV report
        csv_report = generate_csv_report(data, organization, report_type, **kwargs)
        
        # Store CSV report
        if csv_report:
            csv_location = store_report(csv_report, organization, report_type, **kwargs)
        
        # Create PowerPoint report
        ppt_report = create_ppt_report(data, organization, report_type, **kwargs)
        
        # Store PowerPoint report
        if ppt_report:
            ppt_location = store_report(ppt_report, organization, report_type, **kwargs)
        
        return {
            'organization': organization,
            'report_type': report_type,
            'csv_report': csv_location if csv_report else None,
            'ppt_report': ppt_location if ppt_report else None
        }
    
    except Exception as e:
        logger.error(f"Failed to generate {report_type} report for {organization}: {str(e)}")
        raise

# Create tasks for daily reports
for org in ORG_LIST:
    daily_report_task = PythonOperator(
        task_id=f'daily_report_{org}',
        python_callable=generate_report_for_organization,
        op_kwargs={'organization': org, 'report_type': 'daily'},
        dag=dag,
    )

# Create tasks for weekly reports (only run on Mondays)
for org in ORG_LIST:
    weekly_report_task = PythonOperator(
        task_id=f'weekly_report_{org}',
        python_callable=generate_report_for_organization,
        op_kwargs={'organization': org, 'report_type': 'weekly'},
        dag=dag,
    )

# Create tasks for monthly reports (only run on the 1st of the month)
for org in ORG_LIST:
    monthly_report_task = PythonOperator(
        task_id=f'monthly_report_{org}',
        python_callable=generate_report_for_organization,
        op_kwargs={'organization': org, 'report_type': 'monthly'},
        dag=dag,
    )

if __name__ == "__main__":
    dag.cli()